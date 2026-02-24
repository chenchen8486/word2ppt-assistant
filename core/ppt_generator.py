import os
import re
import importlib.util
from pptx import Presentation
from pptx.util import Inches

class PPTGenerator:
    def __init__(self, mcp_server_path=None):
        self.mcp_server_path = mcp_server_path or os.environ.get("PPT_MCP_SERVER_PATH", "D:/project/python/Office-PowerPoint-MCP-Server-main")

    def generate_pptx(self, structured_text, output_file):
        slides = self._parse_structured_text(structured_text)
        if not slides:
            return False, output_file, "No slides generated from structured content."
        prs = Presentation()
        for slide_data in slides:
            slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text = slide_data.get("title", "")
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.2), Inches(5.2))
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            self._fill_body(body_frame, slide_data.get("body", ""))
            image_path = slide_data.get("image", "")
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(9.2), Inches(1.7), width=Inches(3.8))
        prs.save(output_file)
        return True, output_file, "PPTX 生成完成"

    def _fill_body(self, text_frame, body_text):
        text_frame.clear()
        lines = body_text.splitlines() if body_text else []
        first_set = False
        for line in lines:
            if not line.strip():
                continue
            text = line.strip()
            if text.startswith("-"):
                text = text.lstrip("-").strip()
            if not first_set:
                p = text_frame.paragraphs[0]
                p.text = text
                first_set = True
            else:
                p = text_frame.add_paragraph()
                p.text = text

    def _parse_structured_text(self, structured_text):
        slides = []
        blocks = re.findall(r"\[\[SLIDE\]\]([\s\S]*?)\[\[/SLIDE\]\]", structured_text or "")
        for block in blocks:
            title = ""
            body_lines = []
            image = ""
            mode = ""
            for line in block.splitlines():
                if line.startswith("TITLE:"):
                    title = line[len("TITLE:"):].strip().replace("\\_", "_")
                    mode = ""
                    continue
                if line.startswith("BODY:"):
                    mode = "body"
                    content = line[len("BODY:"):].lstrip().replace("\\_", "_")
                    if content:
                        body_lines.append(content)
                    continue
                if line.startswith("IMAGE:"):
                    image = line[len("IMAGE:"):].strip()
                    mode = ""
                    continue
                if mode == "body":
                    body_lines.append(line.replace("\\_", "_"))
            body = "\n".join(body_lines).strip()
            if not title and body:
                first_line = body.splitlines()[0] if body.splitlines() else ""
                title = first_line[:50]
            if title or body:
                slides.append({"title": title, "body": body, "image": image})
        return slides

    def _optimize_with_mcp(self, pptx_path):
        root = self.mcp_server_path
        if not root or not os.path.exists(root):
            return None
        validation_path = os.path.join(root, "utils", "validation_utils.py")
        if not os.path.exists(validation_path):
            return None
        validation_utils = self._load_module("mcp_validation_utils", validation_path)
        if not validation_utils:
            return None
        try:
            pres = Presentation(pptx_path)
            for slide in pres.slides:
                validation_utils.validate_and_fix_slide(slide, auto_fix=True, min_font_size=10, max_font_size=48)
            optimized_path = self._build_optimized_path(pptx_path)
            pres.save(optimized_path)
            return optimized_path
        except Exception:
            return None

    def _load_module(self, name, path):
        spec = importlib.util.spec_from_file_location(name, path)
        if not spec or not spec.loader:
            return None
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)
        return module

    def _build_optimized_path(self, pptx_path):
        base, ext = os.path.splitext(pptx_path)
        return f"{base}_optimized{ext}"
