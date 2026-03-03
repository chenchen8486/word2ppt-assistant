import json
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any


class PPTXGenerator:
    """
    基于模板驱动的PPTX生成器
    使用python-pptx库，根据预定义的template.pptx模板生成PPTX文件
    """

    def __init__(self):
        # 查找模板文件
        self.template_path = self._find_template()
        if self.template_path is None:
            print("警告: 未找到 template.pptx 模板文件")
            print("请确保 data/template.pptx 文件存在")

    def _find_template(self) -> str:
        """
        查找模板文件，支持多种路径
        """
        from pathlib import Path

        # 尝试多个可能的模板位置
        potential_paths = [
            "data/template.pptx",  # 开发环境
            "template.pptx",       # 根目录
            Path(__file__).parent.parent / "data" / "template.pptx",  # 相对于当前文件
        ]

        # 打包环境检查
        if getattr(sys, 'frozen', False):
            # 打包后，模板可能在 exe 同级目录的 data/ 下
            exe_dir = Path(sys.executable).parent
            potential_paths.append(exe_dir / "data" / "template.pptx")

        for path in potential_paths:
            try:
                if Path(path).exists():
                    print(f"找到模板文件: {Path(path).absolute()}")
                    return str(Path(path).absolute())
            except:
                continue

        return None

    def _get_weighted_length(self, text: str) -> int:
        """
        计算文本的加权长度，考虑换行对版面纵向空间的占用

        Args:
            text: 输入文本

        Returns:
            加权长度（字符数 + 换行数 * 35）
        """
        return sum(35 if char == '\n' else 1 for char in text)

    def _split_text_by_sentences(self, text: str, max_weight: int, first_weight: int = None) -> List[str]:
        """
        按句子边界和权重精确分割文本，引入首块贪心限制

        Args:
            text: 输入文本
            max_weight: 最大权重限制
            first_weight: 首块权重限制（用于贪心填充）

        Returns:
            按句子边界和权重分割后的文本列表
        """
        text_weight = self._get_weighted_length(text)
        if text_weight <= max_weight:
            if first_weight is None or text_weight <= first_weight:
                return [text]

        # 使用正则分割，将换行符也作为句子边界
        parts = re.split(r'([。！？；\n])', text)

        # 将分隔符和文本部分组合起来
        sentences = []
        i = 0
        while i < len(parts):
            if i + 1 < len(parts) and parts[i+1] in ['。', '！', '？', '；', '\n']:
                # 将文本部分和标点符号组合
                sentence = parts[i] + parts[i+1]
                sentences.append(sentence)
                i += 2
            else:
                # 最后一部分（如果没有标点符号）
                if parts[i].strip():  # 如果非空
                    sentences.append(parts[i])
                i += 1

        # 实现贪心逻辑
        chunks = []
        current_chunk = ""
        current_weight = 0
        target_weight = first_weight if first_weight is not None else max_weight
        first_chunk_done = False

        for sentence in sentences:
            sentence_weight = self._get_weighted_length(sentence)

            # 如果当前句子本身就已经超过权重限制，需要强制拆分
            if sentence_weight > max_weight:
                # 先将当前积累的块添加到结果中
                if current_chunk:
                    chunks.append(current_chunk)

                # 对超长句子进行强制拆分
                split_chunks = self._brute_force_split(sentence, max_weight)
                chunks.extend(split_chunks)

                # 重置当前块
                current_chunk = ""
                current_weight = 0

                # 如果这是第一个块且已经添加了内容，恢复使用max_weight
                if not first_chunk_done and chunks:
                    first_chunk_done = True
                    target_weight = max_weight

                continue

            # 检查是否超出当前目标权重
            if current_weight + sentence_weight <= target_weight:
                current_chunk += sentence
                current_weight += sentence_weight
            else:
                # 超出当前目标权重，保存当前块
                if current_chunk:
                    chunks.append(current_chunk)

                # 如果这是第一个块，之后的块使用max_weight
                if not first_chunk_done:
                    first_chunk_done = True
                    target_weight = max_weight

                # 如果当前句子可以直接放入新块
                if sentence_weight <= max_weight:
                    current_chunk = sentence
                    current_weight = sentence_weight
                else:
                    # 如果单个句子超限，强制拆分
                    split_chunks = self._brute_force_split(sentence, max_weight)
                    chunks.extend(split_chunks)

                    current_chunk = ""
                    current_weight = 0

        # 添加最后的块
        if current_chunk:
            chunks.append(current_chunk)

        # 处理没有句子边界的情况
        if not chunks and text:
            # 使用暴力拆分
            chunks = self._brute_force_split(text, max_weight)

        return chunks

    def _brute_force_split(self, text: str, max_weight: int) -> List[str]:
        """
        当句子边界分割无法满足权重限制时，强制按字符数分割

        Args:
            text: 输入文本
            max_weight: 最大权重

        Returns:
            分割后的文本列表
        """
        chunks = []
        start = 0

        while start < len(text):
            # 找到合适的切分点
            end = start + max_weight // 2  # 初始估计
            while end < len(text):
                chunk = text[start:end]
                if self._get_weighted_length(chunk) > max_weight:
                    break
                end += 1

            # 回退到合适的点
            while end > start:
                chunk = text[start:end]
                if self._get_weighted_length(chunk) <= max_weight:
                    chunks.append(chunk)
                    start = end
                    break
                end -= 1

            # 如果连一个字符都无法容纳，至少要放入一个字符
            if start < len(text) and end == start:
                chunks.append(text[start:start+1])
                start += 1

        return chunks

    def _create_context_slides(self, prs, content: str, number: str = ""):
        """
        为context内容创建幻灯片，处理长文本分页

        Args:
            prs: Presentation对象
            content: 要渲染的内容
            number: 题号（可选）
        """
        # 拼接题号和内容
        if number:
            full_content = f"{number} {content}"
        else:
            full_content = content

        # 按句子分割长文本，使用1200权重的饱和阈值（符合14pt下的空间优化）
        content_chunks = self._split_text_by_sentences(full_content, 1200)

        for chunk in content_chunks:
            # 创建新幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[self.context_layout_idx])

            # 寻找合适的文本框来填充内容
            target_shape = None
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == self.context_placeholder_idx:
                        target_shape = shape
                        break

            # 如果找到了目标占位符，填充内容
            if target_shape and target_shape.has_text_frame:
                # 设置段落属性
                text_frame = target_shape.text_frame
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                p.text = chunk
                p.font.size = Pt(14)  # 设置字体大小为14pt
                # 如果段落以阿拉伯数字开头，使用黑体，否则使用微软雅黑
                if re.match(r'^\d+', p.text.strip()):
                    p.font.name = '黑体'
                else:
                    p.font.name = '微软雅黑'  # 强制设置中文字体
                p.alignment = PP_ALIGN.LEFT  # 左对齐

    def _create_question_slides(self, prs, content: str, answer: str, analysis: str, number: str = ""):
        """
        为question内容创建幻灯片，使用单个占位符的富文本渲染，带智能分页

        Args:
            prs: Presentation对象
            content: 题干内容
            answer: 答案内容
            analysis: 解析内容
            number: 题号（可选）
        """
        # 拼接题号和题干内容
        if number:
            full_content = f"{number}. {content}"
        else:
            full_content = content

        # 1. 渲染题干 (Content)
        # 将题干按权重分割成块
        content_chunks = self._split_text_by_sentences(full_content, 800)

        # 循环处理题干chunks
        last_tf = None
        current_page_weight = 0
        for i, chunk in enumerate(content_chunks):
            # 为第一个chunk创建幻灯片，后续的chunk也需要新的幻灯片
            if i == 0:  # 第一个chunk，创建初始幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[self.question_layout_idx])

                # 寻找主文本框 (Body Text placeholder)
                body_shape = None
                max_size = 0

                # 首先尝试找到 Body Text 类型的占位符
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        # 根据模板探测结果，Body/BODY类型的占位符是idx=13，type=2(Body Text) 或 type=1(Text)
                        if shape.placeholder_format.type == 2 or shape.placeholder_format.type == 1:  # Body Text or Text
                            body_shape = shape
                            break
                        # 如果有明确的 idx=13 的占位符，也可以作为候选
                        if hasattr(shape.placeholder_format, 'idx') and shape.placeholder_format.idx == 13:
                            body_shape = shape
                            break
                    # 如果没找到Body Text，尝试找最大的占位符
                    if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                        if shape.width * shape.height > max_size:
                            max_size = shape.width * shape.height
                            body_shape = shape

                # 如果没找到合适的占位符，抛出异常
                if body_shape is None or not body_shape.has_text_frame:
                    raise ValueError("未能找到合适的文本占位符来渲染题目内容")

                # 获取文本框架
                tf = body_shape.text_frame
                last_tf = tf

                # 清空默认文本
                tf.clear()
            else:  # 后续chunks，创建新幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[self.question_layout_idx])

                # 寻找主文本框 (Body Text placeholder)
                body_shape = None
                max_size = 0

                # 首先尝试找到 Body Text 类型的占位符
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        if shape.placeholder_format.type == 2:  # Body Text
                            body_shape = shape
                            break
                    # 如果没找到Body Text，尝试找最大的占位符
                    if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                        if shape.width * shape.height > max_size:
                            max_size = shape.width * shape.height
                            body_shape = shape

                # 如果没找到合适的占位符，抛出异常
                if body_shape is None or not body_shape.has_text_frame:
                    raise ValueError("未能找到合适的文本占位符来渲染题目内容")

                # 获取文本框架
                tf = body_shape.text_frame
                last_tf = tf

                # 清空默认文本
                tf.clear()

            # 写入题干内容
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()  # 获取第一个段落或添加新段落
            p.text = chunk
            p.font.size = Pt(14)  # 设置字体大小为14pt
            p.font.bold = False   # 正常粗细
            # 如果段落以阿拉伯数字开头，使用黑体，否则使用微软雅黑
            if re.match(r'^\d+', chunk.strip()):
                p.font.name = '黑体'
            else:
                p.font.name = '微软雅黑'  # 强制设置中文字体
            p.alignment = PP_ALIGN.LEFT  # 左对齐

            # 更新当前页权重
            current_page_weight = self._get_weighted_length(chunk)

        # 2. 渲染答案 (Answer)
        if answer:
            answer_weight = self._get_weighted_length(answer)

            # 检查当前页是否能容纳答案
            if current_page_weight + answer_weight > 800:
                # 需要翻页
                slide = prs.slides.add_slide(prs.slide_layouts[self.question_layout_idx])

                # 寻找主文本框 (Body Text placeholder)
                body_shape = None
                max_size = 0

                # 首先尝试找到 Body Text 类型的占位符
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        if shape.placeholder_format.type == 2:  # Body Text
                            body_shape = shape
                            break
                    # 如果没找到Body Text，尝试找最大的占位符
                    if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                        if shape.width * shape.height > max_size:
                            max_size = shape.width * shape.height
                            body_shape = shape

                # 如果没找到合适的占位符，抛出异常
                if body_shape is None or not body_shape.has_text_frame:
                    raise ValueError("未能找到合适的文本占位符来渲染题目内容")

                # 获取文本框架
                last_tf = body_shape.text_frame
                # 清空默认文本
                last_tf.clear()

                # 重置当前页权重
                current_page_weight = 0

            # 写入答案
            answer_p = last_tf.add_paragraph()  # 添加新段落
            answer_p.text = f"【答案】{answer}"  # 移除\n前缀以符合零空行原则
            answer_p.font.size = Pt(14)  # 设置字体大小为14pt
            answer_p.font.bold = True    # 加粗
            answer_p.font.color.rgb = RGBColor(220, 53, 69)  # 红色
            # 如果段落以阿拉伯数字开头，使用黑体，否则使用微软雅黑
            if re.match(r'^\d+', answer_p.text.strip()):
                answer_p.font.name = '黑体'
            else:
                answer_p.font.name = '微软雅黑'  # 强制设置中文字体
            answer_p.alignment = PP_ALIGN.LEFT  # 左对齐

            # 更新当前页权重
            current_page_weight += answer_weight

        # 3. 渲染解析 (Analysis)
        if analysis:
            # 计算当前幻灯片的剩余空间
            available_space = self.question_max_chars - current_page_weight

            # 如果剩余空间极小（< 100），直接翻页，给新页完整空间
            if available_space < 100:
                # 需要翻页
                slide = prs.slides.add_slide(prs.slide_layouts[self.question_layout_idx])

                # 寻找主文本框 (Body Text placeholder)
                body_shape = None
                max_size = 0

                # 首先尝试找到 Body Text 类型的占位符
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        if shape.placeholder_format.type == 2:  # Body Text
                            body_shape = shape
                            break
                    # 如果没找到Body Text，尝试找最大的占位符
                    if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                        if shape.width * shape.height > max_size:
                            max_size = shape.width * shape.height
                            body_shape = shape

                # 如果没找到合适的占位符，抛出异常
                if body_shape is None or not body_shape.has_text_frame:
                    raise ValueError("未能找到合适的文本占位符来渲染题目内容")

                # 获取文本框架
                last_tf = body_shape.text_frame
                # 清空默认文本
                last_tf.clear()

                # 重置当前页权重
                current_page_weight = 0
                available_space = self.question_max_chars

            # 将 available_space 传入切分器，做到贪心填充
            analysis_chunks = self._split_text_by_sentences(analysis, max_weight=self.question_max_chars, first_weight=available_space)

            for j, ana_chunk in enumerate(analysis_chunks):
                ana_chunk_weight = self._get_weighted_length(ana_chunk)

                # 检查当前页是否能容纳解析块（第一个chunk会用available_space的剩余空间）
                if j == 0:
                    # 第一个chunk使用可用空间
                    effective_space = self.question_max_chars - current_page_weight
                    if current_page_weight + ana_chunk_weight > self.question_max_chars:
                        # 需要翻页
                        slide = prs.slides.add_slide(prs.slide_layouts[self.question_layout_idx])

                        # 寻找主文本框 (Body Text placeholder)
                        body_shape = None
                        max_size = 0

                        # 首先尝试找到 Body Text 类型的占位符
                        for shape in slide.shapes:
                            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                                if shape.placeholder_format.type == 2:  # Body Text
                                    body_shape = shape
                                    break
                            # 如果没找到Body Text，尝试找最大的占位符
                            if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                                if shape.width * shape.height > max_size:
                                    max_size = shape.width * shape.height
                                    body_shape = shape

                        # 如果没找到合适的占位符，抛出异常
                        if body_shape is None or not body_shape.has_text_frame:
                            raise ValueError("未能找到合适的文本占位符来渲染题目内容")

                        # 获取文本框架
                        last_tf = body_shape.text_frame
                        # 清空默认文本
                        last_tf.clear()

                        # 重置当前页权重
                        current_page_weight = 0
                else:
                    # 后续chunks使用完整空间
                    if current_page_weight + ana_chunk_weight > self.question_max_chars:
                        # 需要翻页
                        slide = prs.slides.add_slide(prs.slide_layouts[self.question_layout_idx])

                        # 寻找主文本框 (Body Text placeholder)
                        body_shape = None
                        max_size = 0

                        # 首先尝试找到 Body Text 类型的占位符
                        for shape in slide.shapes:
                            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                                if shape.placeholder_format.type == 2:  # Body Text
                                    body_shape = shape
                                    break
                            # 如果没找到Body Text，尝试找最大的占位符
                            if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                                if shape.width * shape.height > max_size:
                                    max_size = shape.width * shape.height
                                    body_shape = shape

                        # 如果没找到合适的占位符，抛出异常
                        if body_shape is None or not body_shape.has_text_frame:
                            raise ValueError("未能找到合适的文本占位符来渲染题目内容")

                        # 获取文本框架
                        last_tf = body_shape.text_frame
                        # 清空默认文本
                        last_tf.clear()

                        # 重置当前页权重
                        current_page_weight = 0

                # 写入解析块
                analysis_p = last_tf.add_paragraph()  # 添加新段落
                if j == 0:
                    analysis_p.text = f"【解析】{ana_chunk}"  # 移除\n前缀以符合零空行原则
                else:
                    analysis_p.text = f"[解析续]{ana_chunk}"

                analysis_p.font.size = Pt(14)  # 设置字体大小为14pt
                analysis_p.font.bold = False   # 正常粗细
                analysis_p.font.color.rgb = RGBColor(40, 167, 69)  # 绿色
                # 如果段落以阿拉伯数字开头，使用黑体，否则使用微软雅黑
                if re.match(r'^\d+', analysis_p.text.strip()):
                    analysis_p.font.name = '黑体'
                else:
                    analysis_p.font.name = '微软雅黑'  # 强制设置中文字体
                analysis_p.alignment = PP_ALIGN.LEFT  # 左对齐

                # 更新当前页权重
                current_page_weight += ana_chunk_weight

    def generate(self, json_path: str, template_path: str = None, output_path: str = None, doc_title: str = None) -> bool:
        """
        生成PPTX文件的核心方法

        Args:
            json_path: 输入的JSON文件路径（Phase 2生成的数据）
            template_path: 模板PPTX文件路径（可选，如果不指定则使用默认模板）
            output_path: 输出PPTX文件路径
            doc_title: 文档标题

        Returns:
            是否成功生成
        """
        try:
            # 使用传入的模板路径或默认模板路径
            if template_path is None:
                if self.template_path is None:
                    print("Error: 未指定模板路径且未找到默认模板")
                    return False
                template_path = self.template_path

            # 加载JSON数据
            with open(json_path, 'r', encoding='utf-8-sig') as f:
                extracted_data = json.load(f)

            # 加载模板 - 检查模板文件是否存在
            if not Path(template_path).exists():
                print(f"Error: 模板文件不存在: {template_path}")
                return False

            prs = Presentation(template_path)

            # 检查是否有初始幻灯片，并将其作为封面页
            if len(prs.slides) > 0:
                # 如果存在封面页，设置标题
                first_slide = prs.slides[0]
                if first_slide.shapes.title:
                    if doc_title:
                        first_slide.shapes.title.text = doc_title
                else:
                    # 尝试找到其他可能的标题形状
                    for shape in first_slide.shapes:
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            # 检查是否是标题形状（通过索引或其他标识判断）
                            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                                if shape.placeholder_format.type == 0:  # Title placeholder
                                    shape.text_frame.text = doc_title or ""
                                    break
            else:
                # 如果没有初始幻灯片，添加封面页
                if doc_title:
                    title_slide_layout = prs.slide_layouts[0]  # 通常为标题版式
                    title_slide = prs.slides.add_slide(title_slide_layout)
                    if title_slide.shapes.title:
                        title_slide.shapes.title.text = doc_title

            # 遍历所有条目
            for i, item_or_list in enumerate(extracted_data):
                if isinstance(item_or_list, list):
                    # 如果是列表，遍历其中的每个项目
                    for item in item_or_list:
                        self._process_item(prs, item, i == 0)  # 传递首项标识
                elif isinstance(item_or_list, dict):
                    # 如果是字典，直接处理，标记是否为首项
                    is_first_item = (i == 0)
                    self._process_item(prs, item_or_list, is_first_item)

            # 保存PPTX文件
            prs.save(output_path)
            print(f"Successfully saved PPTX to {output_path}")
            return True

        except Exception as e:
            import traceback
            error_msg = f"Error generating PPTX: {str(e)}\n{traceback.format_exc()}"
            print(error_msg)
            # 尝试写入崩溃日志
            try:
                from pathlib import Path
                # 尝试多个可能的日志位置
                log_paths = [
                    Path("crash_log.txt"),
                    Path("data/02_temp_build/crash_log.txt"),
                    Path("dist/Word2PPT-Assistant/crash_log.txt")
                ]
                for log_path in log_paths:
                    try:
                        log_path.parent.mkdir(parents=True, exist_ok=True)
                        with open(log_path, 'w', encoding='utf-8') as f:
                            f.write(error_msg)
                        break
                    except:
                        continue
            except:
                pass
            return False

    def _process_item(self, prs, item: Dict[str, Any], is_first_item: bool = False):
        """
        处理单个项目

        Args:
            prs: Presentation对象
            item: 数据项
            is_first_item: 是否是列表中的第一项
        """
        if not isinstance(item, dict):
            return

        # 检查是否是错误项
        if 'error' in item:
            # 跳过错误项
            print(f"Skipping error item: {item.get('error', 'Unknown error')}")
            return

        # 检查是否缺少必要的字段
        if 'type' not in item:
            return

        item_type = item.get('type', '').lower()

        if item_type == 'context':
            content = item.get('content', '')
            number = item.get('number', '')  # 提取题号

            # 泛化剔除冗余标题（首项智能屏蔽）：如果第一项是context类型且number为空字符串，
            # 极可能是LLM提取的试卷头部标题信息，直接跳过
            if is_first_item and number == "":
                print(f"Skipping redundant header context: '{content[:50]}...' (detected as header)")
                return

            # 智能降噪：如果内容长度小于阈值，跳过
            if len(content.strip()) < self.noise_threshold:
                print(f"Skipping short context: '{content[:50]}...' (length: {len(content)})")
                return

            # 文本净化：压缩连续换行符
            content = re.sub(r'\n\s*\n', '\n', content).strip()

            # 创建context幻灯片
            self._create_context_slides(prs, content, number)

        elif item_type == 'question':
            content = item.get('content', '')
            answer = item.get('answer', '')
            analysis = item.get('analysis', '')
            number = item.get('number', '')  # 提取题号

            # 文本净化：压缩连续换行符
            content = re.sub(r'\n\s*\n', '\n', content).strip()
            if answer:
                answer = re.sub(r'\n\s*\n', '\n', answer).strip()
            if analysis:
                analysis = re.sub(r'\n\s*\n', '\n', analysis).strip()

            # 创建question幻灯片
            self._create_question_slides(prs, content, answer, analysis, number)


if __name__ == "__main__":
    # 测试代码
    generator = PPTXGenerator()
    success = generator.generate(
        json_path="data/02_temp_build/test_extracted.json",
        template_path="data/template.pptx",
        output_path="data/03_output_pptx/test_output.pptx",
        doc_title="测试试卷标题"  # 示例文档标题
    )
    print(f"Generation {'successful' if success else 'failed'}")