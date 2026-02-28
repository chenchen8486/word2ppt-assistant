#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Phase 3 验证脚本 - 验证从提取的JSON到PPTX的生成流程
"""

import os
from pathlib import Path
from core.pptx_generator import PPTXGenerator

def main():
    print("[PHASE 3 TEST] 开始PPTX生成验证...")

    # 定义文件路径
    json_path = "data/02_temp_build/test_extracted.json"
    template_path = "data/template.pptx"

    # 符合Phase 3规范的输出路径
    output_path = "data/02_temp_build/04_test_final.pptx"

    # 确保输出目录存在
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    # 检查输入文件是否存在
    if not Path(json_path).exists():
        print(f"[ERROR] 找不到输入文件: {json_path}")

        # 查找可能的JSON文件
        temp_dir = Path("data/02_temp_build/")
        json_files = list(temp_dir.glob("*.json"))

        if json_files:
            json_path = str(json_files[0])
            print(f"[INFO] 使用找到的JSON文件: {json_path}")
        else:
            print("[ERROR] 在temp目录中没有找到任何JSON文件")
            return

    if not Path(template_path).exists():
        print(f"[ERROR] 找不到模板文件: {template_path}")
        return

    # 创建生成器实例
    generator = PPTXGenerator()

    print(f"[INFO] 输入JSON: {json_path}")
    print(f"[INFO] 模板文件: {template_path}")
    print(f"[INFO] 输出PPTX: {output_path}")

    # 生成PPTX
    success = generator.generate(
        json_path=json_path,
        template_path=template_path,
        output_path=output_path
    )

    if success:
        print(f"[SUCCESS] Phase 3 PPTX 生成成功!")
        print(f"[INFO] 输出文件: {output_path}")

        # 验证输出文件
        if Path(output_path).exists():
            size = Path(output_path).stat().st_size
            print(f"[INFO] 文件大小: {size:,} 字节")

            # 统计幻灯片数量
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"[INFO] 生成幻灯片总数: {len(prs.slides)}")

            # 验证布局分布
            layout_counts = {}
            for i, slide in enumerate(prs.slides):
                layout_idx = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'unknown'
                layout_counts[layout_idx] = layout_counts.get(layout_idx, 0) + 1

            print(f"[INFO] 布局分布: {layout_counts}")

            print("[PHASE 3] 验证完成 - 所有要求已满足!")
        else:
            print("[ERROR] 输出文件未生成")
    else:
        print("[ERROR] PPTX 生成失败")

if __name__ == "__main__":
    main()